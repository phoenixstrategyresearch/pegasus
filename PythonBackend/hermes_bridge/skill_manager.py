"""
Skill manager - persistent skill creation and retrieval.
Ported from Hermes Agent's skill system.

Skills are stored as directories under ~/Documents/pegasus_data/skills/
Each skill has a SKILL.md with YAML frontmatter and markdown body.
"""

import os
import re

SKILLS_DIR = os.path.join(
    os.environ.get("PEGASUS_DATA_DIR", os.path.expanduser("~/Documents/pegasus_data")),
    "skills"
)

VALID_NAME = re.compile(r"^[a-z0-9][a-z0-9._-]{0,63}$")


DEFAULT_SKILLS = {
    "context-optimizer": {
        "description": "Intelligent context window management. Compacts, summarizes, and prioritizes information to maximize effective use of the model context window.",
        "content": """# Context Optimizer

## Purpose
Manage the context window intelligently to prevent overflow while preserving critical information.

## When to Use
- Long conversations that approach context limits
- Processing large documents or web scrapes
- Multi-step research tasks with lots of tool results

## Strategy: Adaptive Compaction

### 1. Semantic Compaction
When multiple messages cover the same topic, merge them into a single summary:
- Identify overlapping information across messages
- Keep the most specific/recent version of each fact
- Discard redundant phrasing

### 2. Temporal Compaction
For older parts of the conversation:
- Summarize exchanges older than 10 turns into bullet points
- Preserve: user preferences, key decisions, action items
- Discard: pleasantries, failed attempts, verbose tool output

### 3. Extractive Compaction
For large tool results (web scrapes, file contents):
- Extract only sentences relevant to the user's query
- Strip boilerplate: navigation, footers, ads, disclaimers
- Keep: data, facts, quotes, code blocks

### 4. Priority Scoring
Score each piece of context:
- **HIGH**: User instructions, current task, recent decisions
- **MEDIUM**: Tool results relevant to current task, established preferences
- **LOW**: Greetings, old tangents, superseded information
- **DROP**: Raw HTML, duplicate content, failed tool attempts

## Implementation
When context is getting long:
1. Identify the current active task/question
2. Score all context by relevance to that task
3. Summarize LOW items into 1-2 sentences
4. Drop items with no relevance
5. Keep HIGH and MEDIUM items intact
""",
    },
    "financial-analyst": {
        "description": "Financial analysis and research workflows for market research, equity research, comparable companies, precedent transactions, and DCF valuation.",
        "content": """# Financial Analyst

## Overview

Deliver industry-standard market research, equity research, comps, precedents, and DCF valuation with explicit assumptions, source citations, and polished outputs in Excel, PowerPoint, or Markdown.

## Quick Intake

- Confirm company name, ticker, industry, geography, currency, and units (USD millions, etc.).
- Confirm scope (DCF, comps, precedents, market sizing, competitive strategy, equity research memo) and output formats.
- Gather data sources (CSV/Excel/SQL/API/web) and ask for access details or files.
- Ask for missing inputs; if the user says proceed, use standard assumptions and list them explicitly.

## Core Workflows

### DCF Model

- Build a 5-year forecast, WACC via CAPM, and terminal value via Gordon Growth.
- Produce a valuation summary, sensitivity table, and assumptions table.

### Comparable Companies

- Select peers with key operating metrics and trading multiples.
- Build a comps table. Call out outliers and justify adjustments or exclusions.

### Precedent Transactions

- Build a precedents table with deal values, premiums, and transaction multiples.
- Note deal context (strategic vs. financial buyer, control premium, timing).

### Market Research

- Provide TAM/SAM/SOM with both top-down and bottom-up triangulation when possible.
- Cite sources and state any conversions or extrapolations.

### Competitive Strategy and Key Competitors

- Identify direct, adjacent, and substitute competitors with brief rationale.
- Assess differentiation, moats, pricing power, and distribution advantages.

### Equity Research Memo

- Present a clear thesis, valuation summary, catalysts, and risks.
- Keep conclusions linked to data and comps/DCF outputs.

## Output Standards

- Provide an assumptions table for every model or valuation.
- Include units, currency, and as-of dates on every table.
- Keep formulas transparent in Excel outputs and document key drivers in Markdown.
""",
    },
    "office-documents": {
        "description": "Create and edit Word (.docx), Excel (.xlsx), and PowerPoint (.pptx) files programmatically using python-docx, openpyxl, and python-pptx.",
        "content": """# Office Document Specialist

## Overview

Create and manipulate professional Office documents on-device. Uses python-docx, openpyxl, and python-pptx.

## Setup

Before first use, install dependencies:
- pip_install("python-docx") for Word documents
- pip_install("openpyxl") for Excel spreadsheets
- pip_install("python-pptx") for PowerPoint presentations

## Word Documents (.docx)

Use python-docx via python_exec:
```python
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH

doc = Document()
doc.add_heading('Report Title', 0)
doc.add_paragraph('Body text here.')

# Add a table
table = doc.add_table(rows=3, cols=3, style='Table Grid')
table.cell(0, 0).text = 'Header'

# Formatting
p = doc.add_paragraph()
run = p.add_run('Bold text')
run.bold = True
run.font.size = Pt(14)

doc.save('report.docx')
```

## Excel Spreadsheets (.xlsx)

Use openpyxl or the excel_read tool for reading. For creation:
```python
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border

wb = Workbook()
ws = wb.active
ws.title = "Data"

# Headers with styling
headers = ['Name', 'Value', 'Date']
for col, h in enumerate(headers, 1):
    cell = ws.cell(row=1, column=col, value=h)
    cell.font = Font(bold=True)
    cell.fill = PatternFill('solid', fgColor='4472C4')
    cell.font = Font(bold=True, color='FFFFFF')

# Formulas
ws['B10'] = '=SUM(B2:B9)'

# Auto-width columns
for col in ws.columns:
    max_len = max(len(str(c.value or '')) for c in col)
    ws.column_dimensions[col[0].column_letter].width = max_len + 2

wb.save('output.xlsx')
```

## PowerPoint Presentations (.pptx)

Use python-pptx:
```python
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Presentation Title"
slide.placeholders[1].text = "Subtitle"

# Content slide with bullets
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Key Points"
tf = slide.placeholders[1].text_frame
tf.text = "First point"
p = tf.add_paragraph()
p.text = "Second point"
p.level = 1

# Add table
rows, cols = 4, 3
table = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(3)).table
table.cell(0, 0).text = "Header"

prs.save('presentation.pptx')
```

## Best Practices

- Always install the required package before using it
- Save files to the workspace root (they'll be accessible to the user)
- For large Excel files, use excel_read tool instead of openpyxl directly
- Apply professional formatting: consistent fonts, colors, alignment
- Use table styles in Word docs for clean output
- Add slide numbers and footers to PowerPoint decks
""",
    },
    "self-improving": {
        "description": "Learn from corrections and self-reflection. Store insights in memory, identify patterns, and compound knowledge across sessions.",
        "content": """# Self-Improving Agent

## Purpose
Continuously improve by learning from user corrections, self-reflection, and pattern recognition. Store lessons in persistent memory.

## Learning Triggers

### Log a correction when:
- User says "No, that's not right" or corrects your output
- User says "I prefer X, not Y" or "Always do X"
- User says "Remember that..." or "Stop doing X"
- You identify a mistake in your own completed work

### Ignore (don't log):
- One-time instructions ("this time, do X")
- Context-specific statements that won't repeat
- Hypothetical questions

## Self-Reflection Protocol

After completing significant work, run this checklist:

1. **Assess**: Did the output match the user's intent? Were there misunderstandings?
2. **Identify**: What could be improved? Was there unnecessary work?
3. **Pattern Check**: Is this a repeating situation? Have I made this mistake before?
4. **Log**: If a pattern is found, save to memory:
   - CONTEXT: [what type of task]
   - LESSON: [what to do differently]

## Memory Management

### What to store:
- User preferences (formatting, tone, tools they like)
- Common mistakes and their fixes
- Project-specific patterns and conventions
- Workflow shortcuts that worked well

### What NOT to store:
- Credentials, API keys, or sensitive data
- One-time task details
- Information that changes frequently

### Memory hygiene:
- Before adding, check if a similar entry exists - update instead of duplicate
- Remove entries the user explicitly contradicts
- Keep entries concise: one clear sentence per lesson

## Commands
- "What have you learned?" -> Read and summarize memory contents
- "Show my patterns" -> Display stored preferences and patterns
- "Forget X" -> Remove specific entries from memory
""",
    },
}


class SkillManager:
    def __init__(self):
        os.makedirs(SKILLS_DIR, exist_ok=True)
        self._install_defaults()

    def _install_defaults(self):
        """Install default skills if they don't exist yet."""
        for name, info in DEFAULT_SKILLS.items():
            skill_dir = os.path.join(SKILLS_DIR, name)
            skill_file = os.path.join(skill_dir, "SKILL.md")
            if not os.path.exists(skill_file):
                os.makedirs(skill_dir, exist_ok=True)
                frontmatter = "name: " + name + "\ndescription: " + info["description"] + "\n"
                content = "---\n" + frontmatter + "---\n\n" + info["content"]
                with open(skill_file, "w", encoding="utf-8") as f:
                    f.write(content)
                print("[Skills] Installed default skill: " + name)

    def _skill_dir(self, name: str, category: str = "") -> str:
        if category:
            return os.path.join(SKILLS_DIR, category, name)
        return os.path.join(SKILLS_DIR, name)

    def list_skills(self) -> list[dict]:
        """List all skills with name, description, category."""
        result = []
        for root, dirs, files in os.walk(SKILLS_DIR):
            if "SKILL.md" not in files:
                continue
            skill_path = os.path.join(root, "SKILL.md")
            meta = self._parse_frontmatter(skill_path)
            rel = os.path.relpath(root, SKILLS_DIR)
            parts = rel.split(os.sep)
            category = parts[0] if len(parts) > 1 else ""
            result.append({
                "name": meta.get("name", parts[-1]),
                "description": meta.get("description", ""),
                "category": category,
            })
        return result

    def view_skill(self, name: str, file_path: str = "") -> dict:
        """View a skill's content or a specific file within it."""
        # Search for the skill
        skill_dir = self._find_skill(name)
        if skill_dir is None:
            return {"error": f"Skill not found: {name}"}

        if file_path:
            full = os.path.join(skill_dir, file_path)
            if not os.path.isfile(full):
                return {"error": f"File not found: {file_path}"}
            with open(full, "r", errors="replace") as f:
                return {"name": name, "file": file_path, "content": f.read(50_000)}

        # Return SKILL.md + file listing
        skill_md = os.path.join(skill_dir, "SKILL.md")
        with open(skill_md, "r", encoding="utf-8") as f:
            content = f.read()

        files = []
        for root, dirs, fnames in os.walk(skill_dir):
            for fn in fnames:
                if fn == "SKILL.md":
                    continue
                rel = os.path.relpath(os.path.join(root, fn), skill_dir)
                files.append(rel)

        return {"name": name, "content": content, "files": files}

    def create_skill(self, name: str, description: str, content: str, category: str = "") -> dict:
        """Create a new skill."""
        if not VALID_NAME.match(name):
            return {"error": f"Invalid skill name: {name}. Use lowercase letters, numbers, hyphens."}
        if len(description) > 1024:
            return {"error": "Description must be <= 1024 characters"}

        skill_dir = self._skill_dir(name, category)
        if os.path.exists(skill_dir):
            return {"error": f"Skill already exists: {name}"}

        os.makedirs(skill_dir, exist_ok=True)

        frontmatter = f"name: {name}\ndescription: {description}\n"
        skill_md = f"---\n{frontmatter}---\n\n{content}"

        with open(os.path.join(skill_dir, "SKILL.md"), "w", encoding="utf-8") as f:
            f.write(skill_md)

        return {"status": "created", "name": name, "path": skill_dir}

    def delete_skill(self, name: str) -> dict:
        """Delete a skill."""
        skill_dir = self._find_skill(name)
        if skill_dir is None:
            return {"error": f"Skill not found: {name}"}
        import shutil
        shutil.rmtree(skill_dir)
        return {"status": "deleted", "name": name}

    def _find_skill(self, name: str) -> str | None:
        """Find a skill directory by name."""
        for root, dirs, files in os.walk(SKILLS_DIR):
            if os.path.basename(root) == name and "SKILL.md" in files:
                return root
        return None

    def _parse_frontmatter(self, path: str) -> dict:
        """Parse simple key: value frontmatter from a SKILL.md file."""
        try:
            with open(path, "r", encoding="utf-8") as f:
                content = f.read()
            if not content.startswith("---"):
                return {}
            end = content.index("---", 3)
            result = {}
            for line in content[3:end].strip().splitlines():
                if ": " in line:
                    key, value = line.split(": ", 1)
                    result[key.strip()] = value.strip()
            return result
        except Exception:
            return {}


# Global singleton
skills = SkillManager()
